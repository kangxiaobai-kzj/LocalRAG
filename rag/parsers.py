# ============================================
# 项目：LocalRAG · 可定制本地 RAG 智能体
# 开发者：kangxiaobai-kzj
# 开发时间：2026-08-14
# ============================================

# rag/parsers.py
# PDF 解析引擎分层：pypdf 快路径 / Tesseract OCR / MinerU（可选）
# parse_pdf() 按"文字版优先 → 扫描件 OCR → 可选 MinerU"分级调度。
import glob
import os
import re
import shutil
import subprocess
import tempfile
from typing import Dict, List, Tuple

from config import PARSER_ENABLE_MINERU
from utils.logger import get_logger

logger = get_logger("parsers")

# ============================================
# 【路径配置】项目内嵌依赖（免系统 PATH）
# ============================================
BASE_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))

# Poppler 路径（用于 PDF 渲染）
POPPLER_PATH = os.path.join(BASE_DIR, "bin", "poppler-26.02.0", "Library", "bin")
if os.path.exists(POPPLER_PATH):
    logger.info("检测到内嵌 Poppler：%s", POPPLER_PATH)
else:
    POPPLER_PATH = None
    logger.warning("未检测到内嵌 Poppler，将尝试系统 PATH")

# Tesseract 路径（用于 OCR 识别）
TESSERACT_PATH = os.path.join(BASE_DIR, "bin", "Tesseract-OCR", "tesseract.exe")
if os.path.exists(TESSERACT_PATH):
    logger.info("检测到内嵌 Tesseract：%s", TESSERACT_PATH)
    try:
        import pytesseract
        pytesseract.pytesseract.tesseract_cmd = TESSERACT_PATH
        logger.info("pytesseract 已配置为使用内嵌 Tesseract")
    except ImportError:
        logger.warning("pytesseract 未安装，请执行：pip install pytesseract")
else:
    TESSERACT_PATH = None
    logger.warning("未检测到内嵌 Tesseract，OCR 将失败")

# OCR 依赖可用性
try:
    from pdf2image import convert_from_path
    PDF2IMAGE_AVAILABLE = True
except ImportError:
    PDF2IMAGE_AVAILABLE = False
    logger.warning("pdf2image 未安装，OCR 将不可用。请安装：pip install pdf2image")

try:
    import pytesseract
    PYTESSERACT_AVAILABLE = True
except ImportError:
    PYTESSERACT_AVAILABLE = False
    logger.warning("pytesseract 未安装，OCR 将不可用。请安装：pip install pytesseract")

# 判定为文字版 PDF 的最低字符数（低于此值视为疑似扫描件）
TEXT_PDF_MIN_CHARS = 50

# ============================================
# 支持的文件类型（上传 / 构建 / 统计共用）
# ============================================
SUPPORTED_EXTENSIONS = {".pdf", ".txt", ".md", ".docx", ".xlsx", ".pptx"}


def is_supported_file(name: str) -> bool:
    """按扩展名判断文件是否可解析入库（忽略大小写）。"""
    return os.path.splitext(name or "")[1].lower() in SUPPORTED_EXTENSIONS


class BaseParser:
    """解析引擎基类。extract 返回 {页码: [文本]}，文本保留换行结构（供标题感知切片使用）。"""

    name = "base"

    def extract(self, file_path: str) -> Dict[int, List[str]]:
        raise NotImplementedError


class PypdfParser(BaseParser):
    """文字版 PDF 快速提取（保留换行）。"""

    name = "pypdf"

    def extract(self, file_path: str) -> Dict[int, List[str]]:
        page_texts: Dict[int, List[str]] = {}
        try:
            from pypdf import PdfReader
            reader = PdfReader(file_path)
            for page_num, page in enumerate(reader.pages, start=1):
                raw = page.extract_text()
                if raw:
                    clean = re.sub(r"[ \t]+", " ", raw).strip()
                    if clean:
                        page_texts[page_num] = [clean]
        except Exception as e:
            logger.warning("pypdf 快速提取失败: %s", e)
        return page_texts


class OcrParser(BaseParser):
    """扫描件 OCR：pdf2image 转图 + pytesseract 识别（中文简体 + 英文）。"""

    name = "tesseract"

    def extract(self, file_path: str) -> Dict[int, List[str]]:
        page_texts: Dict[int, List[str]] = {}
        if not PDF2IMAGE_AVAILABLE or not PYTESSERACT_AVAILABLE:
            logger.warning("OCR 依赖库缺失，跳过")
            return page_texts
        if TESSERACT_PATH is None:
            logger.warning("Tesseract 未找到，跳过 OCR")
            return page_texts
        try:
            images = convert_from_path(file_path, poppler_path=POPPLER_PATH, dpi=200)
            logger.info("共 %d 页，开始 OCR 识别...", len(images))
            for page_num, img in enumerate(images, start=1):
                text = pytesseract.image_to_string(img, lang='chi_sim+eng', config='--psm 6')
                if text.strip():
                    page_texts[page_num] = [text]
                img.close()
            if page_texts:
                logger.info("OCR 成功提取 %d 页内容", len(page_texts))
            else:
                logger.warning("OCR 未提取到任何文本（可能图片清晰度不足）")
        except Exception as e:
            logger.error("OCR 识别失败: %s", e)
            if "poppler" in str(e).lower():
                logger.warning("请检查 Poppler 路径是否正确")
            elif "tesseract" in str(e).lower():
                logger.warning("请检查 Tesseract 路径是否正确，并确保语言包 chi_sim 已安装")
        return page_texts


class MineruParser(BaseParser):
    """可选解析引擎：MinerU（复杂版面/表格/公式）。默认关闭（config.PARSER_ENABLE_MINERU）。"""

    name = "mineru"

    def extract(self, file_path: str) -> Dict[int, List[str]]:
        if not shutil.which("mineru"):
            logger.warning("未检测到 mineru 命令，跳过 MinerU 解析（可选：pip install mineru）")
            return {}
        out_dir = tempfile.mkdtemp(prefix="mineru_")
        try:
            # MinerU CLI 参数（-p 输入 -o 输出）；具体用法以所装版本为准
            subprocess.run(["mineru", "-p", file_path, "-o", out_dir],
                           capture_output=True, timeout=600)
            page_texts: Dict[int, List[str]] = {}
            md_files = sorted(glob.glob(os.path.join(out_dir, "**", "*.md"), recursive=True))
            for i, md in enumerate(md_files, start=1):
                with open(md, "r", encoding="utf-8") as f:
                    content = f.read().strip()
                if content:
                    page_texts[i] = [content]
            if page_texts:
                logger.info("MinerU 解析成功（%d 页）", len(page_texts))
            return page_texts
        except Exception as e:
            logger.error("MinerU 解析失败: %s", e)
            return {}


def parse_pdf(file_path: str) -> Tuple[Dict[int, List[str]], str, int]:
    """
    分级调度：pypdf 快路径 → 疑似扫描件走 OCR → 仍为空且启用时走 MinerU。
    返回 (page_texts, parser_name, char_count)。
    """
    page_texts = PypdfParser().extract(file_path)
    char_count = sum(len(t) for texts in page_texts.values() for t in texts)

    if char_count > TEXT_PDF_MIN_CHARS:
        logger.info("识别为文字版 PDF（%d 字符），快速处理", char_count)
        return page_texts, "pypdf", char_count

    logger.info("疑似扫描件（仅 %d 字符），启动 OCR...", char_count)
    ocr_texts = OcrParser().extract(file_path)
    if ocr_texts:
        ocr_count = sum(len(t) for texts in ocr_texts.values() for t in texts)
        return ocr_texts, "tesseract", ocr_count

    if PARSER_ENABLE_MINERU:
        mineru_texts = MineruParser().extract(file_path)
        if mineru_texts:
            mineru_count = sum(len(t) for texts in mineru_texts.values() for t in texts)
            return mineru_texts, "mineru", mineru_count

    # 全部失败：返回 pypdf 结果（可能为空），由调用方判定失败
    return page_texts, "pypdf", char_count


# ============================================
# 文本类解析：TXT / MD / DOCX / XLSX
# 统一返回 {页码: [文本]}：TXT/MD/DOCX 无分页概念记作第 1 页；XLSX 以工作表序号作"页码"。
# ============================================
def _read_text_file(file_path: str) -> str:
    """按编码优先级读取纯文本文件（utf-8 → gbk → utf-16 → latin-1 兜底）。"""
    for enc in ("utf-8", "gbk", "utf-16", "latin-1"):
        try:
            with open(file_path, "r", encoding=enc) as f:
                return f.read()
        except (UnicodeDecodeError, UnicodeError):
            continue
    return ""


def parse_txt(file_path: str) -> Dict[int, List[str]]:
    """纯文本（TXT）：整文件作为第 1 页返回。"""
    text = _read_text_file(file_path).strip()
    if not text:
        return {}
    return {1: [text]}


def parse_md(file_path: str) -> Dict[int, List[str]]:
    """Markdown（MD）：按纯文本读取，保留标题结构供标题感知切片使用。"""
    text = _read_text_file(file_path).strip()
    if not text:
        return {}
    return {1: [text]}


def parse_docx(file_path: str) -> Dict[int, List[str]]:
    """Word（DOCX）：段落文本 + 表格内容（单元格以 | 连接）合并为第 1 页。"""
    try:
        from docx import Document
    except ImportError:
        logger.warning("python-docx 未安装，无法解析 DOCX。请执行：pip install python-docx")
        return {}
    try:
        doc = Document(file_path)
        parts = []
        for para in doc.paragraphs:
            t = para.text.strip()
            if t:
                parts.append(t)
        for table in doc.tables:
            for row in table.rows:
                cells = [c.text.strip() for c in row.cells]
                line = " | ".join(c for c in cells if c)
                if line:
                    parts.append(line)
        text = "\n".join(parts).strip()
        return {1: [text]} if text else {}
    except Exception as e:
        logger.error("DOCX 解析失败: %s", e)
        return {}


def parse_xlsx(file_path: str) -> Dict[int, List[str]]:
    """Excel（XLSX）：每个工作表为一"页"，行以 | 连接为文本行（data_only 取公式结果）。"""
    try:
        from openpyxl import load_workbook
    except ImportError:
        logger.warning("openpyxl 未安装，无法解析 XLSX。请执行：pip install openpyxl")
        return {}
    try:
        wb = load_workbook(file_path, read_only=True, data_only=True)
        page_texts: Dict[int, List[str]] = {}
        for idx, ws in enumerate(wb.worksheets, start=1):
            lines = []
            for row in ws.iter_rows(values_only=True):
                cells = [str(c).strip() for c in row if c is not None and str(c).strip()]
                if cells:
                    lines.append(" | ".join(cells))
            text = "\n".join(lines).strip()
            if text:
                page_texts[idx] = [text]
        wb.close()
        return page_texts
    except Exception as e:
        logger.error("XLSX 解析失败: %s", e)
        return {}


def parse_pptx(file_path: str) -> Dict[int, List[str]]:
    """PowerPoint（PPTX）：每张幻灯片为一"页"，提取各形状中的文本框与表格文本。
    仅提取文字框/表格中的可读文本；嵌入的图片内容无法直接提取（建议图片版面转 PDF + OCR）。
    若某页仅有图片/无文本，则跳过该页。"""
    try:
        from pptx import Presentation
    except ImportError:
        logger.warning("python-pptx 未安装，无法解析 PPTX。请执行：pip install python-pptx")
        return {}
    try:
        prs = Presentation(file_path)
        page_texts: Dict[int, List[str]] = {}
        for idx, slide in enumerate(prs.slides, start=1):
            lines = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = "".join(run.text for run in para.runs).strip()
                        if t:
                            lines.append(t)
                elif shape.shape_type == 19:  # MSO_SHAPE_TYPE.TABLE
                    try:
                        for row in shape.table.rows:
                            cells = [c.text.strip() for c in row.cells if c.text.strip()]
                            if cells:
                                lines.append(" | ".join(cells))
                    except Exception:
                        pass
            text = "\n".join(lines).strip()
            if text:
                page_texts[idx] = [text]
        return page_texts
    except Exception as e:
        logger.error("PPTX 解析失败: %s", e)
        return {}


# 文本类解析器注册表：扩展名 -> 解析函数
_TEXT_PARSERS = {
    ".txt": parse_txt,
    ".md": parse_md,
    ".docx": parse_docx,
    ".xlsx": parse_xlsx,
    ".pptx": parse_pptx,
}


def parse_document(file_path: str) -> Tuple[Dict[int, List[str]], str, int]:
    """
    按扩展名分发解析（PDF 走原有分级调度，其余走对应文本解析器）。
    返回 (page_texts, parser_name, char_count)。
    """
    ext = os.path.splitext(file_path)[1].lower()
    if ext == ".pdf":
        return parse_pdf(file_path)
    parser = _TEXT_PARSERS.get(ext)
    if parser is None:
        return {}, "unsupported", 0
    page_texts = parser(file_path)
    char_count = sum(len(t) for texts in page_texts.values() for t in texts)
    return page_texts, parser.__name__.replace("parse_", ""), char_count
