# ============================================
# 项目：LocalRAG · 可定制本地 RAG 智能体
# 开发者：kangxiaobai-kzj
# 开发时间：2026-08-14
# ============================================

# 测试：多格式解析器（PDF 原有测试在别处，这里覆盖 TXT/MD/DOCX/XLSX 与分发逻辑）
import os
import tempfile

from docx import Document
from openpyxl import Workbook

from rag.parsers import (
    is_supported_file,
    parse_document,
    parse_md,
    parse_txt,
    SUPPORTED_EXTENSIONS,
)


def test_supported_extensions():
    assert is_supported_file("a.pdf") and is_supported_file("A.TXT")
    assert is_supported_file("b.md") and is_supported_file("c.docx")
    assert is_supported_file("d.xlsx") and is_supported_file("e.pptx")
    assert not is_supported_file("f.doc")
    assert ".pptx" in SUPPORTED_EXTENSIONS


def test_parse_txt():
    with tempfile.NamedTemporaryFile("w", suffix=".txt", encoding="utf-8", delete=False) as f:
        f.write("第一行内容。\n第二行。")
        path = f.name
    try:
        texts, name, n = parse_document(path)
        assert name == "txt" and n > 0
        assert 1 in texts and "第一行内容" in texts[1][0]
    finally:
        os.remove(path)


def test_parse_md():
    with tempfile.NamedTemporaryFile("w", suffix=".md", encoding="utf-8", delete=False) as f:
        f.write("# 标题\n\n正文段落。")
        path = f.name
    try:
        texts, name, n = parse_document(path)
        assert name == "md" and n > 0
        assert "# 标题" in texts[1][0]  # 保留 Markdown 结构供标题感知切片
    finally:
        os.remove(path)


def test_parse_docx():
    with tempfile.TemporaryDirectory() as d:
        path = os.path.join(d, "c.docx")
        doc = Document()
        doc.add_paragraph("Word 段落内容。")
        doc.add_paragraph("第二段。")
        doc.save(path)
        texts, name, n = parse_document(path)
        assert name == "docx" and n > 0
        assert "Word 段落内容" in texts[1][0]


def test_parse_xlsx_sheet_as_page():
    with tempfile.TemporaryDirectory() as d:
        path = os.path.join(d, "d.xlsx")
        wb = Workbook()
        ws = wb.active
        ws.title = "参数表"
        ws.append(["产品", "价格"])
        ws.append(["高精定位终端", 1200])
        wb.save(path)
        texts, name, n = parse_document(path)
        assert name == "xlsx" and n > 0
        assert 1 in texts  # 工作表序号作为"页码"
        assert "高精定位终端" in texts[1][0] and "1200" in texts[1][0]


def test_parse_pptx():
    from pptx import Presentation
    with tempfile.TemporaryDirectory() as d:
        path = os.path.join(d, "e.pptx")
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        box = slide.shapes.add_textbox(0, 0, 100, 50)
        box.text = "PPT 演示内容"
        prs.save(path)
        texts, name, n = parse_document(path)
        assert name == "pptx" and n > 0
        assert 1 in texts and "PPT 演示内容" in texts[1][0]


def test_parse_unsupported_extension():
    with tempfile.NamedTemporaryFile("w", suffix=".doc", delete=False) as f:
        f.write("x")
        path = f.name
    try:
        texts, name, n = parse_document(path)
        assert name == "unsupported" and not texts and n == 0
    finally:
        os.remove(path)


def test_parse_txt_direct():
    with tempfile.NamedTemporaryFile("w", suffix=".txt", encoding="utf-8", delete=False) as f:
        f.write("直接调用测试")
        path = f.name
    try:
        assert 1 in parse_txt(path)
    finally:
        os.remove(path)


def test_parse_md_direct():
    with tempfile.NamedTemporaryFile("w", suffix=".md", encoding="utf-8", delete=False) as f:
        f.write("直接调用测试")
        path = f.name
    try:
        assert 1 in parse_md(path)
    finally:
        os.remove(path)
